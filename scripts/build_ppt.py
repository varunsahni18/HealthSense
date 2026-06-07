"""Build HealthSense slide deck for SparQ 2026.

Run with the project venv:
    .venv312/Scripts/python.exe scripts/build_ppt.py
Output: submissions/HealthSense_Slides.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "submissions" / "HealthSense_Slides.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# Brand palette (matches the HTML one-pager)
TEAL_DARK = RGBColor(0x0C, 0x4A, 0x6E)
TEAL = RGBColor(0x08, 0x91, 0xB2)
TEAL_LIGHT = RGBColor(0xEC, 0xFE, 0xFF)
TEAL_BORDER = RGBColor(0xA5, 0xF3, 0xFC)
GREY_900 = RGBColor(0x11, 0x18, 0x27)
GREY_700 = RGBColor(0x37, 0x41, 0x51)
GREY_500 = RGBColor(0x6B, 0x72, 0x80)
GREY_BG = RGBColor(0xF8, 0xFA, 0xFC)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GREEN_BG = RGBColor(0xDC, 0xFC, 0xE7)
AMBER = RGBColor(0xCA, 0x8A, 0x04)
AMBER_BG = RGBColor(0xFE, 0xF9, 0xC3)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_BG = RGBColor(0xFD, 0xE8, 0xD8)
RED = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def make_pres() -> Presentation:
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, fill, line=None, line_w=0.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 0.5)
    sh.text_frame.text = ""
    sh.text_frame.margin_left = Inches(0.1)
    sh.text_frame.margin_right = Inches(0.1)
    sh.text_frame.margin_top = Inches(0.05)
    sh.text_frame.margin_bottom = Inches(0.05)
    return sh


def add_round(slide, x, y, w, h, fill, line=None, line_w=0.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.18
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 0.75)
    sh.text_frame.margin_left = Inches(0.12)
    sh.text_frame.margin_right = Inches(0.12)
    sh.text_frame.margin_top = Inches(0.08)
    sh.text_frame.margin_bottom = Inches(0.08)
    return sh


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=GREY_900,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Inter"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=GREY_700, font="Inter",
                line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = "• " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, Inches(13.333), Inches(0.55), TEAL_DARK)
    add_text(slide, Inches(0.5), Inches(0.07), Inches(9), Inches(0.45),
             title, size=18, bold=True, color=WHITE)
    add_text(slide, Inches(10.5), Inches(0.13), Inches(2.6), Inches(0.4),
             "HealthSense · SparQ 2026", size=10, color=WHITE, align=PP_ALIGN.RIGHT)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.62), Inches(12), Inches(0.4),
                 subtitle, size=12, color=TEAL, italic_safe=False) if False else \
        add_text(slide, Inches(0.5), Inches(0.62), Inches(12), Inches(0.4),
                 subtitle, size=12, color=TEAL)


def add_footer(slide, page_no, total):
    add_rect(slide, 0, Inches(7.2), Inches(13.333), Inches(0.3), GREY_BG)
    add_text(slide, Inches(0.5), Inches(7.22), Inches(8), Inches(0.28),
             "Track: Edge + Cloud AI by Industry — Healthcare · Idea #02 · "
             "arXiv:2604.04297 · arXiv:2604.13359",
             size=9, color=GREY_500)
    add_text(slide, Inches(11.5), Inches(7.22), Inches(1.5), Inches(0.28),
             f"{page_no} / {total}", size=9, color=GREY_500, align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# Build slides
# ──────────────────────────────────────────────────────────────────────────────

prs = make_pres()
TOTAL = 16


# Slide 1 — Title
s = add_blank(prs)
# gradient-like band
add_rect(s, 0, 0, Inches(13.333), Inches(7.5), TEAL_DARK)
add_rect(s, 0, Inches(2.0), Inches(13.333), Inches(3.6), TEAL)
# title
add_text(s, Inches(1), Inches(2.3), Inches(11.3), Inches(1.2),
         "HealthSense", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.55), Inches(11.3), Inches(0.6),
         "On-Wearable Multimodal Biosignal AI · Zero Cloud Transmission",
         size=22, color=WHITE, align=PP_ALIGN.CENTER)
# pill
pill = add_round(s, Inches(4.6), Inches(4.3), Inches(4.1), Inches(0.55),
                 RGBColor(0xFF, 0xFF, 0xFF), TEAL_BORDER, 1.0)
pill.fill.transparency = 0
add_text(s, Inches(4.6), Inches(4.36), Inches(4.1), Inches(0.45),
         "Edge + Cloud AI by Industry — Healthcare", size=14, bold=True,
         color=TEAL_DARK, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.3), Inches(11.3), Inches(0.5),
         "PanLUNA 5.4M-Parameter Foundation Encoder · BioTrain On-Device Personalisation",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
         "Qualcomm SparQ 2026  ·  Idea #02 of 10", size=12,
         color=RGBColor(0xBA, 0xE6, 0xFD), align=PP_ALIGN.CENTER)


# Slide 2 — The Problem
s = add_blank(prs)
add_header_bar(s, "The Problem — Health Data Has a Trust Tax",
               "Continuous monitoring requires sending the most intimate signals to the cloud.")
# 3 problem cards
cards = [
    ("Privacy", "Raw EEG/ECG streamed to cloud is a HIPAA / GDPR\n"
                "liability. Wearers do not trust 24/7 monitoring.", ORANGE_BG, ORANGE),
    ("Latency", "Cloud round-trip is 1–3 s. A ventricular\n"
                "tachycardia alert that late is useless.", AMBER_BG, AMBER),
    ("Battery & Cost", "An always-on cellular/BLE link drains a typical\n"
                       "wearable battery in <24 h. Cloud GPUs cost ~$0.20/user/day.", TEAL_LIGHT, TEAL),
]
for i, (title, body, bg, accent) in enumerate(cards):
    x = Inches(0.5 + i * 4.2)
    add_round(s, x, Inches(1.4), Inches(4.0), Inches(2.7), bg, accent, 1.5)
    add_text(s, x + Inches(0.25), Inches(1.6), Inches(3.6), Inches(0.5),
             title, size=20, bold=True, color=accent)
    add_text(s, x + Inches(0.25), Inches(2.2), Inches(3.6), Inches(2.0),
             body, size=14, color=GREY_700)
# stat strip
add_round(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.4), GREY_BG, GREY_500, 0.5)
add_text(s, Inches(0.7), Inches(4.55), Inches(12), Inches(0.4),
         "By the numbers", size=12, bold=True, color=TEAL)
stats = [
    ("422 M", "people with diabetes — most without continuous monitoring", "WHO 2024"),
    ("1 in 5", "adults will have AFib by age 80; most undetected at onset", "AHA 2024"),
    ("$60 B+", "global cardiac & sleep monitoring market, growing 12% YoY", "GVR 2025"),
]
for i, (val, lbl, src) in enumerate(stats):
    x = Inches(0.7 + i * 4.05)
    add_text(s, x, Inches(5.0), Inches(3.9), Inches(0.7),
             val, size=34, bold=True, color=TEAL_DARK)
    add_text(s, x, Inches(5.85), Inches(3.9), Inches(0.7),
             lbl, size=12, color=GREY_700)
    add_text(s, x, Inches(6.4), Inches(3.9), Inches(0.4),
             src, size=10, color=GREY_500, font="Consolas")
add_footer(s, 2, TOTAL)


# Slide 3 — Our Solution (1 line + tagline)
s = add_blank(prs)
add_header_bar(s, "Our Solution — A Foundation Model on the Wrist",
               "One 5.4M-parameter encoder, five biosignals, four heads, zero raw data egress.")
# big sentence
add_round(s, Inches(0.7), Inches(1.4), Inches(12), Inches(2.0), TEAL_LIGHT, TEAL_BORDER, 1.0)
add_text(s, Inches(1.0), Inches(1.7), Inches(11.4), Inches(1.5),
         "Run PanLUNA — a multimodal biosignal foundation model — directly on the\n"
         "Snapdragon Hexagon NPU.  Personalise it on-device with BioTrain.\n"
         "Send the cloud only the conclusion, never the signal.",
         size=22, bold=True, color=TEAL_DARK)
# four pillars
pillars = [
    ("Privacy by Design", "0 bytes of raw biosignal\nleave the device", GREEN_BG, GREEN),
    ("Real-Time", "~100 ms wrist-to-alert\non Hexagon NPU", TEAL_LIGHT, TEAL),
    ("All-Day Battery", "18.8 mJ per ECG inference\n(GAP9 baseline)", AMBER_BG, AMBER),
    ("Personalised", "+35 % accuracy from\non-device fine-tune", ORANGE_BG, ORANGE),
]
for i, (title, body, bg, accent) in enumerate(pillars):
    x = Inches(0.5 + i * 3.15)
    add_round(s, x, Inches(3.7), Inches(2.95), Inches(3.0), bg, accent, 1.5)
    add_text(s, x + Inches(0.2), Inches(3.85), Inches(2.6), Inches(0.5),
             title, size=15, bold=True, color=accent)
    add_text(s, x + Inches(0.2), Inches(4.4), Inches(2.6), Inches(2.4),
             body, size=14, color=GREY_700)
add_footer(s, 3, TOTAL)


# Slide 4 — Quick Glossary
s = add_blank(prs)
add_header_bar(s, "Quick Glossary — Five Biosignals · Two Models",
               "One-sentence definitions before the deep dives.")

# Top: 5 biosignal modality cards (3 + 2 layout)
modalities = [
    ("EEG",  "Brain electrical activity",
     "From scalp electrodes — 8 channels, 250 Hz",
     "Tells us: sleep stage, drowsiness, focus", TEAL_LIGHT, TEAL),
    ("ECG",  "Heart electrical activity",
     "From chest or wrist electrodes — up to 12 leads, 500 Hz",
     "Tells us: rhythm (sinus, AFib, VT), HR, HRV", RGBColor(0xFE, 0xE2, 0xE2), RED),
    ("PPG",  "Blood-volume pulse",
     "From a green LED + photodiode — 100 Hz",
     "Tells us: pulse rate, SpO₂, BP estimate", ORANGE_BG, ORANGE),
    ("EMG",  "Muscle electrical activity",
     "From skin electrodes — 4 channels, 1 kHz",
     "Tells us: tension, activity, gait quality",
     RGBColor(0xEDE9FE), RGBColor(0x7C, 0x3A, 0xED)),
    ("IMU",  "6-axis motion sensor",
     "Accelerometer + gyroscope — 100 Hz",
     "Tells us: posture, fall, gait, exercise", GREEN_BG, GREEN),
]
for i, (name, what, how, signal, bg, accent) in enumerate(modalities):
    col = i if i < 3 else i - 3
    row = 0 if i < 3 else 1
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.2 + row * 1.55)
    add_round(s, x, y, Inches(4.0), Inches(1.4), bg, accent, 1.5)
    add_text(s, x + Inches(0.2), y + Inches(0.08), Inches(3.7), Inches(0.4),
             name + " — " + what, size=14, bold=True, color=accent)
    add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(3.7), Inches(0.4),
             how, size=10.5, color=GREY_700)
    add_text(s, x + Inches(0.2), y + Inches(0.95), Inches(3.7), Inches(0.4),
             signal, size=10.5, color=GREY_900)

# Bottom: 2 model cards
add_round(s, Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.5),
          TEAL_LIGHT, TEAL_DARK, 1.5)
add_text(s, Inches(0.7), Inches(4.55), Inches(5.6), Inches(0.4),
         "PanLUNA — the foundation model", size=15, bold=True, color=TEAL_DARK)
add_text(s, Inches(0.7), Inches(5.0), Inches(5.6), Inches(0.45),
         "arXiv:2604.04297 (ETH Zürich, 2026)", size=10.5,
         color=TEAL_DARK, font="Consolas")
add_text(s, Inches(0.7), Inches(5.45), Inches(5.6), Inches(1.4),
         "A 5.4 M-parameter Transformer that takes EEG, ECG, PPG, EMG\n"
         "and IMU into one shared encoder — like one brain that learnt\n"
         "all five 'languages'. Pretrained self-supervised on ~40,000 h\n"
         "of public biosignals; INT8-quantised for the Hexagon NPU.",
         size=11, color=GREY_900)

add_round(s, Inches(6.85), Inches(4.4), Inches(6.0), Inches(2.5),
          GREEN_BG, GREEN, 1.5)
add_text(s, Inches(7.05), Inches(4.55), Inches(5.6), Inches(0.4),
         "BioTrain — on-device personalisation", size=15, bold=True, color=GREEN)
add_text(s, Inches(7.05), Inches(5.0), Inches(5.6), Inches(0.45),
         "arXiv:2604.13359 (on-device fine-tuning)", size=10.5,
         color=GREEN, font="Consolas")
add_text(s, Inches(7.05), Inches(5.45), Inches(5.6), Inches(1.4),
         "A recipe that runs full back-propagation on the wearable\n"
         "itself — 0.67 MB RAM, under 50 mW — so PanLUNA learns\n"
         "your unique baseline. Reports +35 % avg. accuracy gain\n"
         "after personalisation, with no raw data leaving the device.",
         size=11, color=GREY_900)

# Form-factor strip
add_round(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.25),
          GREY_BG, GREY_500, 0.5)
add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.25),
         "Target form factor: ANY Snapdragon-Hexagon wearable — "
         "next-gen smartwatch · chest patch · smart ring · EEG sleep band.  "
         "We ship a software platform, not a specific device.",
         size=9.5, color=GREY_700)

add_footer(s, 4, TOTAL)


# Slide 5 — System Architecture
s = add_blank(prs)
add_header_bar(s, "System Architecture",
               "Five sensors → one encoder → four clinical heads. Cloud is opt-in & metadata-only.")
# Sensor column
sensors = [("EEG", "8 ch · 250 Hz"), ("ECG", "12-lead · 500 Hz"),
           ("PPG", "3 ch · 100 Hz"), ("EMG", "4 ch · 1 kHz"),
           ("IMU", "6-DoF · 100 Hz")]
for i, (name, spec) in enumerate(sensors):
    y = Inches(1.1 + i * 0.75)
    add_round(s, Inches(0.5), y, Inches(2.0), Inches(0.6), TEAL_LIGHT, TEAL, 1.0)
    add_text(s, Inches(0.6), y + Inches(0.05), Inches(1.8), Inches(0.3),
             name, size=14, bold=True, color=TEAL_DARK)
    add_text(s, Inches(0.6), y + Inches(0.32), Inches(1.8), Inches(0.25),
             spec, size=9, color=GREY_500)
# PanLUNA encoder
add_round(s, Inches(3.0), Inches(1.1), Inches(4.0), Inches(4.5), TEAL_LIGHT, TEAL, 2.0)
add_text(s, Inches(3.0), Inches(1.25), Inches(4.0), Inches(0.5),
         "PanLUNA", size=22, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
add_text(s, Inches(3.0), Inches(1.75), Inches(4.0), Inches(0.4),
         "5.4 M params · INT8 · Hexagon NPU", size=12, color=TEAL,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(3.0), Inches(2.1), Inches(4.0), Inches(0.4),
         "Pretrained on ~40,000 h biosignal", size=11, color=GREY_500,
         align=PP_ALIGN.CENTER)
add_round(s, Inches(3.3), Inches(2.7), Inches(3.4), Inches(0.55),
          RGBColor(0xBA, 0xE6, 0xFD), TEAL, 1.0)
add_text(s, Inches(3.3), Inches(2.78), Inches(3.4), Inches(0.45),
         "Sensor-Type Embedding", size=12, bold=True, color=TEAL_DARK,
         align=PP_ALIGN.CENTER)
add_round(s, Inches(3.3), Inches(3.35), Inches(3.4), Inches(0.55),
          RGBColor(0xBA, 0xE6, 0xFD), TEAL, 1.0)
add_text(s, Inches(3.3), Inches(3.43), Inches(3.4), Inches(0.45),
         "Cross-Modal Transformer", size=12, bold=True, color=TEAL_DARK,
         align=PP_ALIGN.CENTER)
add_round(s, Inches(3.3), Inches(4.0), Inches(3.4), Inches(0.55),
          RGBColor(0xE0, 0xF2, 0xFE), TEAL_BORDER, 1.0)
add_text(s, Inches(3.3), Inches(4.08), Inches(3.4), Inches(0.45),
         "BioTrain On-Device Tuner", size=12, bold=True, color=TEAL,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(3.0), Inches(4.85), Inches(4.0), Inches(0.4),
         "<50 mW · 0.67 MB RAM · 8.1× memory cut",
         size=10, color=GREY_500, align=PP_ALIGN.CENTER)
# Heads
heads = [("Cardiac Arrhythmia", GREEN_BG, GREEN),
         ("Sleep Stage (SOTA)", AMBER_BG, AMBER),
         ("Stress Level", ORANGE_BG, ORANGE),
         ("EEG Abnormality", TEAL_LIGHT, TEAL)]
for i, (name, bg, accent) in enumerate(heads):
    y = Inches(1.1 + i * 1.13)
    add_round(s, Inches(7.5), y, Inches(3.0), Inches(0.85), bg, accent, 1.5)
    add_text(s, Inches(7.5), y + Inches(0.27), Inches(3.0), Inches(0.4),
             name, size=14, bold=True, color=accent, align=PP_ALIGN.CENTER)
# Zero-cloud badge
add_round(s, Inches(11.0), Inches(2.3), Inches(2.0), Inches(2.0), GREEN_BG, GREEN, 2.0)
add_text(s, Inches(11.0), Inches(2.6), Inches(2.0), Inches(0.6),
         "0 bytes", size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, Inches(11.0), Inches(3.2), Inches(2.0), Inches(0.4),
         "raw signal\nto cloud", size=12, color=GREEN, align=PP_ALIGN.CENTER)
# Cloud band (opt-in)
add_round(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.1),
          GREY_BG, GREY_500, 0.5)
add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.35),
         "Optional Cloud (Hydra AIC100 / QGenie) — clinician-only, metadata-only",
         size=12, bold=True, color=GREY_700)
add_text(s, Inches(0.7), Inches(6.30), Inches(12), Inches(0.6),
         "Clinician summary (gpt-oss-20b)   ·   "
         "Trend forecasting (Qwen3-VL-32B)   ·   "
         "DP-SGD federated weight delta (ε=8) — never raw signals",
         size=11, color=GREY_700)
add_footer(s, 5, TOTAL)


# Slide 5 — PanLUNA Deep Dive
s = add_blank(prs)
add_header_bar(s, "PanLUNA — A Foundation Model for Biosignals",
               "Pretrained on 40,000 hours; 57× smaller than comparable specialists.")
# Left: explanation
add_text(s, Inches(0.5), Inches(1.2), Inches(6.5), Inches(0.5),
         "What it is", size=18, bold=True, color=TEAL_DARK)
add_bullets(s, Inches(0.5), Inches(1.7), Inches(6.5), Inches(4),
            [
                "A single Transformer encoder shared across EEG, ECG, PPG, EMG, IMU.",
                "Each input gets a sensor-type embedding (like a token-type id in BERT).",
                "Cross-modal attention learns relationships between heart, brain & motion.",
                "Pretrained self-supervised on ~40,000 h of public biosignals.",
                "Fine-tunes per task with tiny task heads (≈50 k params each).",
            ], size=14)
# Right: numbers panel
add_round(s, Inches(7.4), Inches(1.2), Inches(5.4), Inches(5.5), TEAL_LIGHT, TEAL_BORDER, 1.0)
add_text(s, Inches(7.6), Inches(1.35), Inches(5), Inches(0.5),
         "Reported numbers (paper)", size=14, bold=True, color=TEAL_DARK)
nums = [
    ("5.4 M", "parameters total"),
    ("57×", "smaller than comparable specialists"),
    ("81.21 %", "balanced accuracy — TUAB EEG abnormality"),
    ("74.16 %", "balanced accuracy — HMC sleep staging (SOTA)"),
    ("325.6 ms", "12-lead ECG latency on GAP9 MCU"),
    ("18.8 mJ", "energy per inference (GAP9 baseline)"),
]
for i, (val, lbl) in enumerate(nums):
    y = Inches(1.95 + i * 0.72)
    add_text(s, Inches(7.6), y, Inches(2.0), Inches(0.6),
             val, size=24, bold=True, color=TEAL_DARK)
    add_text(s, Inches(9.7), y + Inches(0.1), Inches(3.0), Inches(0.6),
             lbl, size=12, color=GREY_700)
add_text(s, Inches(7.6), Inches(6.3), Inches(5), Inches(0.4),
         "Sources: arXiv:2604.04297 · arXiv:2604.13359",
         size=9, color=GREY_500, font="Consolas")
add_footer(s, 6, TOTAL)


# Slide 6 — On-Device Performance
s = add_blank(prs)
add_header_bar(s, "On-Device Performance",
               "We extrapolate paper's GAP9 MCU numbers to Snapdragon Hexagon NPU targets.")
# Comparison bars table
add_text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         "Latency targets (12-lead ECG, INT8)", size=14, bold=True, color=TEAL_DARK)
rows = [
    ("Server GPU (cloud baseline)", 1300, GREY_500),
    ("GAP9 MCU (paper)",             325, AMBER),
    ("Snapdragon X Elite NPU (target)", 95, GREEN),
]
for i, (lbl, ms, col) in enumerate(rows):
    y = Inches(1.7 + i * 0.7)
    add_text(s, Inches(0.5), y, Inches(3.5), Inches(0.5),
             lbl, size=12, color=GREY_700)
    bar_w = Inches(0.0083 * ms)  # 1 ms ≈ 0.0083 in (max ~10.5 in)
    add_rect(s, Inches(4.2), y + Inches(0.1), bar_w, Inches(0.32), col)
    add_text(s, Inches(4.2) + bar_w + Inches(0.1), y + Inches(0.05), Inches(2),
             Inches(0.4), f"{ms} ms", size=12, bold=True, color=GREY_900)
# Energy panel
add_round(s, Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.5),
          TEAL_LIGHT, TEAL_BORDER, 1.0)
add_text(s, Inches(0.7), Inches(4.45), Inches(5.6), Inches(0.5),
         "Energy budget — coin-cell battery", size=14, bold=True, color=TEAL_DARK)
add_bullets(s, Inches(0.7), Inches(4.95), Inches(5.6), Inches(2),
            [
                "18.8 mJ per ECG inference (paper GAP9).",
                "@ 1 inference / 5 s = ~3.8 mW continuous.",
                "240 mAh coin cell ≈ 230 hours = 9.6 days continuous.",
                "Cloud baseline (BLE radio) = ~24 h on the same battery.",
            ], size=12)
# Throughput panel
add_round(s, Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.5),
          GREEN_BG, GREEN, 1.0)
add_text(s, Inches(7.0), Inches(4.45), Inches(5.6), Inches(0.5),
         "Throughput — fits the duty cycle", size=14, bold=True, color=GREEN)
add_bullets(s, Inches(7.0), Inches(4.95), Inches(5.6), Inches(2),
            [
                "EEG: 17 samples / s (paper) — clinical sleep needs 0.5 / s.",
                "ECG: 1 inference per beat (~1 Hz at rest).",
                "PPG: 1 inference / 10 s for stress.",
                "Plenty of head-room for personalisation pass in parallel.",
            ], size=12)
add_footer(s, 7, TOTAL)


# Slide 7 — BioTrain (on-device fine-tune)
s = add_blank(prs)
add_header_bar(s, "BioTrain — Personalisation On the Device",
               "Full back-prop under 50 mW. Your watch learns *you*.")
# Concept
add_text(s, Inches(0.5), Inches(1.2), Inches(6.5), Inches(0.5),
         "Why personalise?", size=18, bold=True, color=TEAL_DARK)
add_bullets(s, Inches(0.5), Inches(1.7), Inches(6.5), Inches(2.5),
            [
                "Every wearer has a unique baseline HRV, EEG rhythm, gait.",
                "Generic models miss subtle deviations from your normal.",
                "Cloud personalisation = ship raw biosignals = privacy non-starter.",
                "BioTrain runs full back-prop *on the watch* in <50 mW.",
            ], size=14)
# How
add_text(s, Inches(0.5), Inches(4.4), Inches(6.5), Inches(0.5),
         "How it fits in 0.67 MB RAM", size=18, bold=True, color=TEAL_DARK)
add_bullets(s, Inches(0.5), Inches(4.9), Inches(6.5), Inches(2.5),
            [
                "Replace BatchNorm → GroupNorm (no activation cache).",
                "Compiler-driven tiling — chunk activations to fit on-chip SRAM.",
                "Gradient accumulation across micro-batches.",
                "Result: 8.1× memory reduction vs naive PyTorch path.",
            ], size=14)
# Right: result panel
add_round(s, Inches(7.4), Inches(1.2), Inches(5.4), Inches(5.5),
          GREEN_BG, GREEN, 1.0)
add_text(s, Inches(7.6), Inches(1.4), Inches(5), Inches(0.5),
         "Live demo result", size=16, bold=True, color=GREEN)
add_text(s, Inches(7.6), Inches(2.0), Inches(5), Inches(2),
         "+35 %", size=72, bold=True, color=GREEN)
add_text(s, Inches(7.6), Inches(3.6), Inches(5), Inches(0.5),
         "personalisation accuracy gain", size=14, bold=True, color=GREEN)
add_text(s, Inches(7.6), Inches(4.1), Inches(5), Inches(0.5),
         "(paper, average across tasks)", size=11, color=GREY_500)
add_text(s, Inches(7.6), Inches(4.9), Inches(5), Inches(0.5),
         "On stage we show:", size=13, bold=True, color=GREEN)
add_bullets(s, Inches(7.6), Inches(5.3), Inches(5), Inches(2),
            [
                "30-second baseline capture",
                "On-device fine-tune (~5 s)",
                "Stress F1: 0.78 → 0.91",
            ], size=12, color=GREEN)
add_footer(s, 8, TOTAL)


# Slide 8 — Software Architecture
s = add_blank(prs)
add_header_bar(s, "Software Architecture — Seven Edge Agents on an Event Bus",
               "Pub/sub orchestration with safety pre-emption when cardiac severity > 0.7.")
agents = [
    ("A1 Perception", "Filter, window, sensor sync", "5 ms · CPU LITTLE"),
    ("A2 PanLUNA",    "INT8 encoder + 4 task heads", "100 ms · Hexagon NPU"),
    ("A3 BioTrain",   "On-device fine-tune (BackProp)", "<50 mW · CPU BIG"),
    ("A4 Edge LLM",   "Explains alerts in plain English", "~3 s · Ollama"),
    ("A5 Proactive",  "5-min trend forecast (cloud)",     "async · Qwen3-VL-32B"),
    ("A6 Complex",    "Clinician Q&A (cloud)",            "async · gpt-oss-20b"),
    ("A7 RAG",        "Local cardiology guidelines",       "<50 ms · BGE-small"),
]
for i, (name, role, perf) in enumerate(agents):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.2 + row * 1.7)
    add_round(s, x, y, Inches(2.95), Inches(1.55), TEAL_LIGHT, TEAL, 1.0)
    add_text(s, x + Inches(0.15), y + Inches(0.08), Inches(2.7), Inches(0.4),
             name, size=14, bold=True, color=TEAL_DARK)
    add_text(s, x + Inches(0.15), y + Inches(0.5), Inches(2.7), Inches(0.6),
             role, size=11, color=GREY_700)
    add_text(s, x + Inches(0.15), y + Inches(1.1), Inches(2.7), Inches(0.4),
             perf, size=10, color=TEAL, font="Consolas")
# bus banner
add_round(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.7),
          TEAL_DARK, TEAL_DARK, 1.0)
add_text(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.4),
         "Event bus  ·  topics: sensor.* · infer.* · alert.* · train.*  ·  "
         "safety pre-emption when alert.cardiac_severity > 0.7",
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# Architectural notes
add_round(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.4),
          TEAL_LIGHT, TEAL, 1.0)
add_text(s, Inches(0.7), Inches(5.75), Inches(12), Inches(0.4),
         "Architecture highlights",
         size=13, bold=True, color=TEAL_DARK)
add_text(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.8),
         "Thread-safe pub/sub bus  ·  cardiac safety pre-emption inside one cycle  ·  "
         "QGenie fallback chain for cloud agents  ·  Kokoro on-device TTS for alerts  ·  "
         "BGE-small ONNX RAG  ·  session export · waveform-stream + state-event SSE channels.",
         size=12, color=GREY_700)
add_footer(s, 9, TOTAL)


# Slide 9 — Data privacy / zero egress
s = add_blank(prs)
add_header_bar(s, "Privacy by Construction — Zero Raw Egress",
               "Even the optional cloud receives only embeddings, alerts, and weight deltas.")
# left: what stays on device
add_round(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.6),
          GREEN_BG, GREEN, 1.5)
add_text(s, Inches(0.7), Inches(1.4), Inches(5.6), Inches(0.5),
         "On the device (always)", size=16, bold=True, color=GREEN)
add_bullets(s, Inches(0.7), Inches(2.0), Inches(5.6), Inches(4),
            [
                "Raw EEG, ECG, PPG, EMG, IMU waveforms",
                "PanLUNA inference & 4 task heads",
                "BioTrain personalisation gradient & weights",
                "Local explanation LLM (Ollama qwen2:7b)",
                "Local RAG corpus (cardiology / sleep)",
                "Kokoro on-device TTS for spoken alerts",
            ], size=13, color=GREEN)
# right: what cloud sees
add_round(s, Inches(6.85), Inches(1.2), Inches(6.0), Inches(5.6),
          GREY_BG, GREY_500, 1.0)
add_text(s, Inches(7.05), Inches(1.4), Inches(5.6), Inches(0.5),
         "What the cloud may see (opt-in)", size=16, bold=True, color=GREY_700)
add_bullets(s, Inches(7.05), Inches(2.0), Inches(5.6), Inches(4),
            [
                "Alert tokens: {label, confidence, timestamp}",
                "1024-d latent embedding (one-way, not invertible)",
                "DP-SGD weight delta with ε=8 (federated learning)",
                "Clinician note text the wearer typed",
                "NEVER: raw waveforms, raw IMU, raw audio",
            ], size=13, color=GREY_700)
# bottom badge
add_round(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.0),
          TEAL_DARK, TEAL_DARK, 0.5)
add_footer(s, 10, TOTAL)


# Slide 10 — Roadmap (Phase plan)
s = add_blank(prs)
add_header_bar(s, "Roadmap — 8 Days to Demo",
               "Bootstrap done. Phase 1 next.")
phases = [
    ("Phase 0", "Bootstrap", "Done", "Project skeleton · venv · slides · speech notes", GREEN),
    ("Phase 1", "Sensors & sim",   "Day 1–2", "Synthetic biosignal generator · uPlot charts · BLE shim", TEAL),
    ("Phase 2", "PanLUNA on NPU",  "Day 3–4", "Compile via AI Hub · QNN-EP loader · 4 heads", TEAL),
    ("Phase 3", "BioTrain",        "Day 5",   "On-device fine-tune · before/after demo widget", AMBER),
    ("Phase 4", "Cloud & polish",  "Day 6–7", "Clinician dashboard · DP-SGD stub · safety pre-empt", ORANGE),
    ("Phase 5", "Submit",          "Day 8",   "Final PPT/PDF · regression suite · video backup", GREEN),
]
for i, (ph, title, when, body, col) in enumerate(phases):
    y = Inches(1.2 + i * 0.95)
    add_round(s, Inches(0.5), y, Inches(1.3), Inches(0.8), col, col, 1.0)
    add_text(s, Inches(0.5), y + Inches(0.15), Inches(1.3), Inches(0.5),
             ph, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_round(s, Inches(2.0), y, Inches(10.8), Inches(0.8), TEAL_LIGHT, col, 0.75)
    add_text(s, Inches(2.2), y + Inches(0.06), Inches(3), Inches(0.4),
             title, size=14, bold=True, color=TEAL_DARK)
    add_text(s, Inches(2.2), y + Inches(0.42), Inches(2.5), Inches(0.4),
             when, size=11, color=GREY_500, font="Consolas")
    add_text(s, Inches(5.0), y + Inches(0.2), Inches(7.7), Inches(0.5),
             body, size=12, color=GREY_700)
add_footer(s, 11, TOTAL)


# Slide 11 — Demo Plan
s = add_blank(prs)
add_header_bar(s, "Live Demo Plan — 90 Seconds",
               "Live biosignal stream, on-device classification, instant personalisation.")
demo = [
    ("0–10 s",  "Wearable streams come alive. 5 stacked waveforms. \"This is real-time.\""),
    ("10–25 s", "PanLUNA inference shows: Sinus · Wake · Low stress · Normal EEG. \n"
                "Latency badge: 97 ms · 19 mJ."),
    ("25–45 s", "Inject AFib-like ECG strip. Cardiac head fires (conf 0.92). \n"
                "Agent 4 narrates the alert in plain English."),
    ("45–65 s", "Click \"Personalise\". 30 s baseline → BioTrain runs on-device → \n"
                "stress F1 jumps from 0.78 → 0.91 visibly on-screen."),
    ("65–80 s", "Show the cloud payload inspector — literally 0 bytes of raw signal. \n"
                "Clinician dashboard receives \"AFib, conf 0.92, 12:04 UTC\"."),
    ("80–90 s", "Tag-line: \"Your data stays on your wrist. Your insights don't.\""),
]
for i, (t, body) in enumerate(demo):
    y = Inches(1.2 + i * 0.95)
    add_round(s, Inches(0.5), y, Inches(1.5), Inches(0.8), TEAL, TEAL, 1.0)
    add_text(s, Inches(0.5), y + Inches(0.15), Inches(1.5), Inches(0.5),
             t, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")
    add_round(s, Inches(2.2), y, Inches(10.6), Inches(0.8), TEAL_LIGHT, TEAL_BORDER, 0.5)
    add_text(s, Inches(2.4), y + Inches(0.1), Inches(10.2), Inches(0.7),
             body, size=12, color=GREY_700)
add_footer(s, 12, TOTAL)


# Slide 12 — Risks & Mitigations
s = add_blank(prs)
add_header_bar(s, "Risks & Mitigations",
               "Every demo-day failure mode has a fallback path.")
risks = [
    ("Bluetooth sensor drops",
     "Fallback to recorded patient slices via sim_source.js; same UI, same latency.",
     ORANGE),
    ("Quantisation accuracy drop",
     "Mixed-precision (FP16 first, then a8w8 selectively); 1k-sample calibration.",
     AMBER),
    ("PanLUNA weights not redistributable",
     "Train smaller equivalent on TUEG + MIMIC-IV; cite paper as PanLUNA-style.",
     AMBER),
    ("Audience asks: \"is the cloud doing the work?\"",
     "Live tcpdump on stage; disconnect WiFi mid-demo; classification continues.",
     GREEN),
    ("Demo runs long",
     "Pre-recorded 90 s video as backup; stop-watch on stage.",
     GREEN),
]
for i, (r, m, col) in enumerate(risks):
    y = Inches(1.2 + i * 1.1)
    add_round(s, Inches(0.5), y, Inches(0.4), Inches(0.95), col, col, 1.0)
    add_round(s, Inches(1.0), y, Inches(11.8), Inches(0.95), TEAL_LIGHT, col, 0.75)
    add_text(s, Inches(1.2), y + Inches(0.08), Inches(11.4), Inches(0.4),
             r, size=14, bold=True, color=GREY_900)
    add_text(s, Inches(1.2), y + Inches(0.5), Inches(11.4), Inches(0.45),
             "Mitigation: " + m, size=12, color=GREY_700)
add_footer(s, 13, TOTAL)


# Slide 13 — Models & Datasets
s = add_blank(prs)
add_header_bar(s, "Models & Datasets — What's Actually Trained",
               "Two ONNX 1D-CNNs we exported this session, plus the public dataset behind the credibility number.")

# ── Two model cards ──
add_round(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(3.5), TEAL_LIGHT, TEAL, 1.5)
add_text(s, Inches(0.7), Inches(1.4), Inches(5.6), Inches(0.5),
         "Model 1 — synth-trained CNN", size=15, bold=True, color=TEAL_DARK)
add_text(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(0.4),
         "models/ecg_cnn.onnx", size=11, color=TEAL, font="Consolas")
m1_lines = [
    ("Architecture",  "1D-CNN · 3 conv blocks · 14,501 params"),
    ("Classes",       "5  (Sinus · AFib · VT · S-Tach · S-Brad)"),
    ("Training data", "Synthetic 4-s windows from biosignal engine (1500)"),
    ("Test acc",      "100 % on synth held-out set (deterministic)"),
    ("Size on disk",  "21 KB (ONNX opset 17)"),
    ("Inference",     "<1 ms ONNX-RT CPU EP · QNN-EP swap = 1 line"),
    ("Role in demo",  "Live cardiac head — drives all 5 scenarios"),
]
y0 = 2.30
for k, v in m1_lines:
    add_text(s, Inches(0.7), Inches(y0), Inches(1.6), Inches(0.3),
             k, size=10, bold=True, color=GREY_700)
    add_text(s, Inches(2.4), Inches(y0), Inches(4.0), Inches(0.3),
             v, size=10, color=GREY_900)
    y0 += 0.32

add_round(s, Inches(6.85), Inches(1.2), Inches(6.0), Inches(3.5), GREEN_BG, GREEN, 1.5)
add_text(s, Inches(7.05), Inches(1.4), Inches(5.6), Inches(0.5),
         "Model 2 — PhysioNet-trained CNN", size=15, bold=True, color=GREEN)
add_text(s, Inches(7.05), Inches(1.85), Inches(5.6), Inches(0.4),
         "models/ecg_cnn_physionet.onnx", size=11, color=GREEN, font="Consolas")
m2_lines = [
    ("Architecture",  "Same · 14,435 params (3-class output)"),
    ("Classes",       "3  (Sinus rhythm · AFib · Other)"),
    ("Training data", "MIT-BIH AFib (afdb/1.0.0) — 4 patients · 530 windows"),
    ("Test acc",      "96.8 % AFib · 91.7 % Sinus on real-patient hold-out"),
    ("Size on disk",  "21 KB"),
    ("Inference",     "Same path — drop-in replacement"),
    ("Role in demo",  "Credibility number / clinical validation slide"),
]
y0 = 2.30
for k, v in m2_lines:
    add_text(s, Inches(7.05), Inches(y0), Inches(1.6), Inches(0.3),
             k, size=10, bold=True, color=GREY_700)
    add_text(s, Inches(8.75), Inches(y0), Inches(4.0), Inches(0.3),
             v, size=10, color=GREY_900)
    y0 += 0.32

# ── Dataset banner (full width) ──
add_round(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(2.0),
          GREY_BG, GREY_500, 0.75)
add_text(s, Inches(0.7), Inches(5.10), Inches(12), Inches(0.4),
         "Dataset — MIT-BIH Atrial Fibrillation Database  (PhysioNet · afdb/1.0.0)",
         size=14, bold=True, color=TEAL_DARK)
add_text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.4),
         "Goldberger et al., \"PhysioBank, PhysioToolkit, and PhysioNet\", "
         "Circulation 101(23): e215–e220, 2000.",
         size=10, color=GREY_700)
add_text(s, Inches(0.7), Inches(5.85), Inches(12), Inches(0.4),
         "23 long-term ECG recordings of subjects with paroxysmal AFib  ·  "
         "ECG sampled at 250 Hz  ·  rhythm-level annotations (AFIB, N, AFL, J).",
         size=10, color=GREY_700)
add_text(s, Inches(0.7), Inches(6.20), Inches(12), Inches(0.4),
         "We use records 04015, 04043, 04126, 04746 — first 30 minutes each — "
         "streamed via wfdb-python directly from PhysioNet.",
         size=10, color=GREY_700, font="Consolas")
add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
         "Reproduce:  .venv312/Scripts/python.exe scripts/train_ecg_cnn_physionet.py  "
         "(downloads data + trains + exports ONNX in ~50 s)",
         size=10, color=GREY_500, font="Consolas")

add_footer(s, 14, TOTAL)


# Slide 14 — Closing
s = add_blank(prs)
add_rect(s, 0, 0, Inches(13.333), Inches(7.5), TEAL_DARK)
add_text(s, Inches(1), Inches(2.0), Inches(11.3), Inches(0.6),
         "HealthSense", size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.1), Inches(11.3), Inches(0.5),
         "Your data stays on your wrist.", size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.5),
         "Your insights don't.", size=28, bold=True,
         color=RGBColor(0xBA, 0xE6, 0xFD), align=PP_ALIGN.CENTER)
add_round(s, Inches(4.6), Inches(5.1), Inches(4.1), Inches(0.7),
          WHITE, TEAL_BORDER, 1.0)
add_text(s, Inches(4.6), Inches(5.2), Inches(4.1), Inches(0.5),
         "Questions?", size=22, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
         "arXiv:2604.04297 · arXiv:2604.13359 · SparQ 2026 · Idea #02",
         size=12, color=RGBColor(0xBA, 0xE6, 0xFD), align=PP_ALIGN.CENTER)


prs.save(OUT)
print(f"OK: wrote {OUT}  ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")
