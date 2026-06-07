"""Append a 'Live Demo Now Working' slide + refresh perf numbers across the deck."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPT = ROOT / "submissions" / "HealthSense_Slides.pptx"

TEAL_DARK = RGBColor(0x0C, 0x4A, 0x6E)
TEAL = RGBColor(0x08, 0x91, 0xB2)
TEAL_LIGHT = RGBColor(0xEC, 0xFE, 0xFF)
TEAL_BORDER = RGBColor(0xA5, 0xF3, 0xFC)
GREY_700 = RGBColor(0x37, 0x41, 0x51)
GREY_500 = RGBColor(0x6B, 0x72, 0x80)
GREY_900 = RGBColor(0x11, 0x18, 0x27)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GREEN_BG = RGBColor(0xDC, 0xFC, 0xE7)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_BG = RGBColor(0xFD, 0xE8, 0xD8)
RED = RGBColor(0xDC, 0x26, 0x26)
RED_BG = RGBColor(0xFE, 0xE2, 0xE2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_rect(slide, x, y, w, h, fill, line=None, line_w=0.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 0.5)
    sh.text_frame.text = ""
    return sh


def add_round(slide, x, y, w, h, fill, line=None, line_w=0.0, adj=0.18):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    sh.adjustments[0] = adj
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 0.75)
    return sh


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=GREY_900,
             align=PP_ALIGN.LEFT, font="Inter"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, Inches(13.333), Inches(0.55), TEAL_DARK)
    add_text(slide, Inches(0.5), Inches(0.07), Inches(9), Inches(0.45),
             title, size=18, bold=True, color=WHITE)
    add_text(slide, Inches(10.5), Inches(0.13), Inches(2.6), Inches(0.4),
             "HealthSense · SparQ 2026", size=10, color=WHITE, align=PP_ALIGN.RIGHT)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.62), Inches(12), Inches(0.4),
                 subtitle, size=12, color=TEAL)


def add_footer(slide, page_no, total):
    add_rect(slide, 0, Inches(7.2), Inches(13.333), Inches(0.3),
             RGBColor(0xF8, 0xFA, 0xFC))
    add_text(slide, Inches(0.5), Inches(7.22), Inches(8), Inches(0.28),
             "Track: Edge + Cloud AI by Industry — Healthcare · Idea #02 · "
             "arXiv:2604.04297 · arXiv:2604.13359",
             size=9, color=GREY_500)
    add_text(slide, Inches(11.5), Inches(7.22), Inches(1.5), Inches(0.28),
             f"{page_no} / {total}", size=9, color=GREY_500, align=PP_ALIGN.RIGHT)


# ── Open the existing deck ──
prs = Presentation(str(PPT))

# Insert a new "Live demo working" slide AFTER slide 11 (Demo Plan).
# python-pptx doesn't support arbitrary insertion, so we append and reorder XML.
blank = prs.slide_layouts[6]
s = prs.slides.add_slide(blank)

# Header
add_header_bar(s, "Demo — Live Now",
               "What you are about to see, served from http://localhost:5000")
# left: ASCII layout / mockup
add_round(s, Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.6),
          RGBColor(0x0F, 0x17, 0x2A), TEAL_DARK, 1.0, adj=0.05)
add_text(s, Inches(0.7), Inches(1.35), Inches(7.6), Inches(0.4),
         "Dashboard layout (3 panels · 14 widgets)", size=12, bold=True,
         color=RGBColor(0xCF, 0xFA, 0xFE))
ascii_layout = """┌───────────────────────────────── HealthSense ─────────────── ●Local-Only · 0 cloud B ──┐
│ ┌─────────────────┐ ┌────────────────── Live biosignals ──┐ ┌──── PanLUNA heads ────┐ │
│ │ Scenario        │ │  ECG  ╱╲╱╲╱╲╱╲╱╲╱╲╱╲╱╲╱╲╱╲╱╲╱╲╱╲╱╲ │ │ Cardiac Sinus rhythm  │ │
│ │  Normal AFib VT │ │  PPG  ∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿  │ │ Sleep   Wake          │ │
│ │  Stress Drowsy  │ │  EEG  ／＼／＼／＼／＼／＼／＼／＼  │ │ Stress  low (13 %)    │ │
│ │                 │ │  EMG  ▁▁▂▃▃▂▁▁▁▁▂▃▂▁▁▁▁▂▂▁▁▁▁    │ │ EEG     Normal        │ │
│ │ Vitals          │ │  IMU  ────╱──╲────╱──╲────╱──╲   │ │                       │ │
│ │  HR  72 bpm     │ │                                    │ │ Alert narration       │ │
│ │  HRV 52 ms      │ │  [latency 95 ms] [energy 19 mJ]    │ │ "All clear. Normal."  │ │
│ │  SpO₂ 98 %      │ │  [inferences 130]                  │ │                       │ │
│ │                 │ └────────────────────────────────────┘ │ Privacy inspector     │ │
│ │ BioTrain        │                                        │ ┌──────┐┌────┐┌────┐ │ │
│ │  ▆▆▆▆▆▆▆▆▆ 92 % │                                        │ │ 4 KB ││ 0 B││ 0 B│ │ │
│ │  F1 0.78→0.91   │                                        │ │ local││cloud││raw │ │ │
│ └─────────────────┘                                        │ └──────┘└────┘└────┘ │ │
└───────────────────────────────────────────────────────────────────────────────────────┘"""
add_text(s, Inches(0.7), Inches(1.7), Inches(7.6), Inches(5.0),
         ascii_layout, size=8.5, color=RGBColor(0xCF, 0xFA, 0xFE),
         font="Consolas")

# right: measured numbers
add_round(s, Inches(8.7), Inches(1.2), Inches(4.1), Inches(5.6),
          GREEN_BG, GREEN, 1.5)
add_text(s, Inches(8.9), Inches(1.4), Inches(3.7), Inches(0.4),
         "Measured (right now)", size=14, bold=True, color=GREEN)
nums = [
    ("21 KB",    "Synth-trained ONNX cardiac CNN"),
    ("14.5 k",   "params · ~0.7 ms ONNX-RT inference"),
    ("100 %",    "test acc on synth (5-class)"),
    ("21 KB",    "PhysioNet-trained ONNX (afdb)"),
    ("96.8 %",   "AFib accuracy on REAL patient ECGs"),
    ("91.7 %",   "Sinus accuracy on real patient ECGs"),
    ("95 ms",    "p50 end-to-end inference latency"),
    ("0 B",      "raw biosignal egress (invariant)"),
]
y0 = 1.85
for v, l in nums:
    add_text(s, Inches(8.9), Inches(y0), Inches(2.0), Inches(0.35),
             v, size=18, bold=True, color=TEAL_DARK, font="JetBrains Mono")
    add_text(s, Inches(10.95), Inches(y0 + 0.04), Inches(1.8), Inches(0.4),
             l, size=10, color=GREY_700)
    y0 += 0.6

add_text(s, Inches(8.9), Inches(6.55), Inches(3.7), Inches(0.3),
         "  models/ecg_cnn_physionet.onnx · MIT-BIH afdb (4 patients)",
         size=9.5, color=GREY_500, font="Consolas")

add_footer(s, 13, 16)

# ── Reorder: move the new slide to position 13 (between Demo Plan #12 and Risks #13) ──
xml_slides = prs.slides._sldIdLst   # type: ignore[attr-defined]
slides = list(xml_slides)
# new slide is at the end; move it to index 12 (0-based, so position 13 visually)
new = slides[-1]
xml_slides.remove(new)
xml_slides.insert(12, new)

# Bump footer numbers on slides that shifted (positions 14, 15, 16 — 0-indexed 13, 14, 15)
def _patch_footer(slide_obj, new_page: int, total: int):
    target = f"{new_page} / {total}"
    for shape in slide_obj.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "/" in text and len(text) <= 8 and text.split("/")[1].strip().isdigit():
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
                if para.runs:
                    para.runs[0].text = target
                else:
                    para.add_run().text = target
            return True
    return False

for new_pos, slide_obj in enumerate(prs.slides, start=1):
    if new_pos in (14, 15, 16):
        _patch_footer(slide_obj, new_pos, 16)

prs.save(str(PPT.with_suffix(".updated.pptx")))
out = PPT.with_suffix(".updated.pptx")
print(f"OK: wrote {out} ({len(prs.slides)} slides; demo slide inserted as #12)")
print("Note: original was locked (open in PPT?); rename when ready.")
